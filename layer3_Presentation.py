# layer3_Presentation.py
"""
PRISM – Layer 3: Whiteboard Presentation Generator
---------------------------------------------------
Reads layer2_notes.txt (Frame-based notes) and creates a PowerPoint presentation.
Each "### Frame" section becomes one slide with bullets underneath.
"""

import os
import re
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor


INPUT_FILE = "layer2_notes.txt"
OUTPUT_FILE = "prism_whiteboard_slides.pptx"


def parse_frames(text: str):
    """
    Split text into frames using headings like '### Frame 1 – ...'
    Returns: list of (frame_title, bullets[])
    """
    frames = []
    # Split on lines starting with '###'
    parts = re.split(r"(?m)^###\s*", text)
    for part in parts:
        part = part.strip()
        if not part:
            continue

        lines = part.splitlines()
        title = lines[0].strip() if lines else "Untitled Frame"
        bullets = []
        for line in lines[1:]:
            line = line.strip()
            if not line:
                continue
            # remove bullet symbols if present
            if line.startswith(("- ", "• ", "* ")):
                bullets.append(line[2:].strip())
            else:
                bullets.append(line)
        frames.append((title, bullets))
    return frames


def create_ppt(frames):
    """
    Builds and saves a PowerPoint file where each frame is a slide.
    """
    prs = Presentation()
    title_font_size = Pt(40)
    bullet_font_size = Pt(22)

    for title, bullets in frames:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        # Title styling
        title_box = slide.shapes.title
        title_box.text = title
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.bold = True
        title_para.font.size = title_font_size
        title_para.font.color.rgb = RGBColor(0, 102, 204)  # PRISM blue

        # Bullets / content
        content = slide.placeholders[1].text_frame
        content.clear()
        if bullets:
            for i, b in enumerate(bullets):
                if i == 0:
                    p = content.paragraphs[0]
                else:
                    p = content.add_paragraph()
                p.text = b
                p.level = 0
                p.font.size = bullet_font_size
                p.font.color.rgb = RGBColor(60, 60, 60)
        else:
            p = content.paragraphs[0]
            p.text = " "

    prs.save(OUTPUT_FILE)
    return OUTPUT_FILE


if __name__ == "__main__":
    if not os.path.exists(INPUT_FILE):
        print(f"⚠️ {INPUT_FILE} not found. Run layer2_Notes.py first.")
        exit(1)

    with open(INPUT_FILE, "r", encoding="utf-8") as f:
        notes_text = f.read()

    frames = parse_frames(notes_text)
    if not frames:
        print("⚠️ No frames found. Make sure notes contain headings like '### Frame 1 – ...'")
        exit(1)

    output_path = create_ppt(frames)
    print(f"\n✅ Presentation created: {output_path}\n")
